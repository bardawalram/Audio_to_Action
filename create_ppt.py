from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Use the PerVoskiteCell.pptx as template to inherit its theme
prs = Presentation('/Users/bardawalparshuram/Downloads/PerVoskiteCell.pptx')

# Delete all existing slides (keep the theme/master)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId_attr = list(prs.slides._sldIdLst[0].attrib.keys())
        for attr in rId_attr:
            if 'id' in attr.lower() and 'r' in attr.lower():
                rId = prs.slides._sldIdLst[0].get(attr)
                break
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

SW = prs.slide_width
SH = prs.slide_height

# Colors for custom shapes
BLUE = RGBColor(0x44, 0x72, 0xC4)       # accent1
ORANGE = RGBColor(0xED, 0x7D, 0x31)     # accent2
GREEN = RGBColor(0x70, 0xAD, 0x47)      # accent6
GOLD = RGBColor(0xFF, 0xC0, 0x00)       # accent4
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5) # accent5
DARK = RGBColor(0x44, 0x54, 0x6A)       # dk2
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x59, 0x59, 0x59)
LIGHT_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
CARD_BG = RGBColor(0xF2, 0xF2, 0xF2)

# Get layout references
layout_title_content = prs.slide_layouts[1]   # Title and Content
layout_section = prs.slide_layouts[2]          # Section Header
layout_blank = prs.slide_layouts[6]            # Blank
layout_comparison = prs.slide_layouts[4]       # Comparison
layout_title_only = prs.slide_layouts[5]       # Title Only


def add_slide(layout=None):
    if layout is None:
        layout = layout_title_content
    return prs.slides.add_slide(layout)


def txt(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def multi_para(slide, left, top, width, items, size=16, color=DARK, bold=False, spacing=Pt(8), bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet:
            p.text = f"• {item}"
        else:
            p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.font.bold = bold
        p.space_after = spacing
    return txBox


def rect(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def bordered_rect(slide, left, top, width, height, fill_color, border_color, border_width=Pt(1.5), radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def line_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
sl = add_slide(layout_title_content)
# Clear default placeholders
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.15), Inches(9), Inches(0.4),
    "Department of Electronics and Communication Engineering", 18, DARK, True, PP_ALIGN.CENTER, "Calibri Light")

txt(sl, Inches(0.5), Inches(0.65), Inches(9), Inches(0.8),
    "Topic : Audio-to-Action: Voice-Driven School ERP System (ReATOA)", 22, DARK, True, PP_ALIGN.CENTER, "Calibri Light")

# Presented by
txBox = sl.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Presented by:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK
p.font.name = "Calibri"
p2 = tf.add_paragraph()
p2.text = "B. Parshuram : 2204206"
p2.font.size = Pt(16)
p2.font.color.rgb = DARK
p2.font.name = "Calibri"
p2.space_before = Pt(8)

# Supervisor
txBox2 = sl.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4.5), Inches(2.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Under the Supervision Of"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK
p.font.name = "Calibri"
lines = [
    "Dr. Girdhar Gopal",
    "(Project Supervisor)",
    "Asst. Professor ECE dept.",
    "NIT-PATNA",
]
for line in lines:
    p2 = tf2.add_paragraph()
    p2.text = f"    {line}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK
    p2.font.name = "Calibri"
    p2.space_before = Pt(4)

# Industry partner
txt(sl, Inches(0.5), Inches(3.6), Inches(9), Inches(0.4),
    "Industry Partner: K12 Techno Services Pvt. Ltd. (SDE)", 14, GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Content / Table of Contents
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Content", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

toc_items = [
    "About K12 — Industry Partner",
    "Introduction",
    "Problem Statement",
    "Objectives",
    "Tech Stack",
    "System Architecture",
    "Voice Processing Pipeline",
    "App Features — Marks & Attendance",
    "Backend & Data Model",
    "Security & Edge Cases",
    "Advantages and Limitations",
    "Future Scope",
    "References",
]

multi_para(sl, Inches(0.7), Inches(0.8), Inches(8), toc_items, 16, DARK, True, Pt(6), bullet=False)


# ============================================================
# SLIDE 3: About K12 — Industry Partner
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "About K12 — Industry Partner", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

bordered_rect(sl, Inches(0.4), Inches(0.8), Inches(9.2), Inches(0.8), CARD_BG, BLUE, Pt(1.5), 0.03)
txt(sl, Inches(0.6), Inches(0.85), Inches(8.8), Inches(0.7),
    "K12 Techno Services is a leading EdTech company focused on delivering high-quality, technology-driven K-12 education solutions across India, partnering with schools to enhance learning outcomes.",
    14, DARK, False, PP_ALIGN.CENTER)

# Three info cards
info = [
    ("Founded", "2009\nBengaluru, India"),
    ("Reach", "500+ partner schools\nacross India"),
    ("Focus", "K-12 curriculum\n& EdTech platforms"),
]
for i, (title, desc) in enumerate(info):
    x = Inches(0.4 + i * 3.15)
    y = Inches(1.85)
    bordered_rect(sl, x, y, Inches(2.95), Inches(1.1), WHITE, BLUE, Pt(1), 0.04)
    txt(sl, x + Inches(0.2), y + Inches(0.1), Inches(2.55), Inches(0.3),
        title, 16, BLUE, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.2), y + Inches(0.45), Inches(2.55), Inches(0.55),
        desc, 12, GRAY, False, PP_ALIGN.CENTER)

# Relevance
txt(sl, Inches(0.5), Inches(3.15), Inches(9), Inches(0.3),
    "RELEVANCE TO REATOA PROJECT", 12, BLUE, True)

relevance = [
    "K12's school management needs directly informed the voice-driven ERP design of ReATOA.",
    "Target users (teachers & staff) align with K12's partner school workforce across India.",
    "K12's technology-first classroom approach validated the voice interface as a viable input method.",
    "The Django REST backend is architected to support future integration with K12's school systems.",
]
for i, r in enumerate(relevance):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 4.8)
    y = Inches(3.5 + row * 0.6)
    bordered_rect(sl, x, y, Inches(4.6), Inches(0.5), CARD_BG, BLUE, Pt(0.75), 0.03)
    txt(sl, x + Inches(0.15), y + Inches(0.07), Inches(4.3), Inches(0.35),
        r, 10, DARK)


# ============================================================
# SLIDE 4: Introduction
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "INTRODUCTION", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

txt(sl, Inches(0.5), Inches(0.8), Inches(9), Inches(0.4),
    "ReATOA (Audio-to-Action)", 24, DARK, True)

intro_text = (
    "ReATOA is a voice-driven school ERP system that enables teachers to manage "
    "marks & attendance through natural spoken commands. It converts speech to text, "
    "extracts intent, confirms with the user, and executes database operations."
)
txt(sl, Inches(0.5), Inches(1.3), Inches(9), Inches(0.8),
    intro_text, 16, DARK)

txt(sl, Inches(0.5), Inches(2.2), Inches(9), Inches(0.4),
    "This project focuses on:", 20, DARK, True)

focus_items = [
    ("1. Voice-first school administration", "    Hands-free marks entry, attendance management via spoken commands."),
    ("2. Custom NLP engine (no ML required)", "    4,000+ lines of regex-based intent & entity extraction from natural speech."),
    ("3. Two-phase commit for safety", "    Preview parsed data, confirm before writing to database. Complete audit trail."),
]

y_pos = 2.7
for title, desc in focus_items:
    txt(sl, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3),
        title, 16, DARK, True)
    txt(sl, Inches(0.5), Inches(y_pos + 0.3), Inches(9), Inches(0.3),
        desc, 14, GRAY)
    y_pos += 0.7


# ============================================================
# SLIDE 5: Problem Statement
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Problem Statement", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

problems = [
    "Teachers spend 30-40% of admin time on repetitive manual data entry for marks & attendance records.",
    "Existing school ERP systems have steep learning curves with complex, multi-step navigation and forms.",
    "No hands-free solution exists for teachers managing students in a classroom setting.",
    "Manual entry leads to frequent errors in student records, with no audit trail to track changes.",
]

for i, prob in enumerate(problems):
    y = Inches(0.9 + i * 1.0)
    # Number
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.5), Inches(0.5))
    c.fill.solid()
    c.fill.fore_color.rgb = BLUE
    c.line.fill.background()
    ctf = c.text_frame
    ctf.paragraphs[0].text = f"0{i+1}"
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    txt(sl, Inches(1.2), y + Inches(0.05), Inches(8.3), Inches(0.5),
        prob, 16, DARK)


# ============================================================
# SLIDE 6: Objectives
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Objectives", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

objectives = [
    "Build a voice-driven ERP system using React + Django + OpenAI Whisper for school administration.",
    "Design a 6-stage voice pipeline: Capture → Transcribe → Upload → Normalize → Extract → Confirm.",
    "Implement custom regex-based NLP (4,000+ LOC) for intent & entity extraction from natural speech.",
    "Develop marks management: subject-wise entry, question-wise scoring & automatic grade calculation.",
    'Enable batch voice operations: "Questions 1,2,3 marks as 4,5,3" in a single command.',
    "Build attendance module with bulk marking & exception handling via voice commands.",
]

for i, obj in enumerate(objectives):
    y = Inches(0.85 + i * 0.7)
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.45), Inches(0.45))
    c.fill.solid()
    c.fill.fore_color.rgb = BLUE
    c.line.fill.background()
    ctf = c.text_frame
    ctf.paragraphs[0].text = f"0{i+1}"
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    txt(sl, Inches(1.1), y + Inches(0.05), Inches(8.4), Inches(0.5),
        obj, 15, DARK)


# ============================================================
# SLIDE 7: Tech Stack
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Tech Stack", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

stacks = [
    ("FRONTEND", "React 18 + Vite 5", "Redux Toolkit · Tailwind CSS\nAxios · React Router 6\nHeroicons · Responsive UI", BLUE),
    ("BACKEND", "Django 5.0 + DRF", "REST APIs · JWT Auth\nPostgreSQL 15\nOpenAI Whisper (STT)", ORANGE),
    ("VOICE / AI", "Custom NLP Engine", "Web Speech API (Live)\nMediaRecorder API\n4,000+ LOC Regex Intent", GREEN),
    ("STATE MGMT", "Redux Toolkit", "Auth Slice · Voice Slice\nUI Slice · Marks State\nLocalStorage persistence", LIGHT_BLUE),
    ("AUDIO", "Dual STT System", "Web Speech API (real-time)\nWhisper (offline fallback)\nMediaRecorder (capture)", GOLD),
    ("DATABASE", "PostgreSQL 15", "10+ models · Foreign Keys\nAudit trail · Seed scripts\nDocker Compose ready", RGBColor(0xA5, 0xA5, 0xA5)),
]

for i, (label, title, desc, color) in enumerate(stacks):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 3.15)
    y = Inches(0.85 + row * 2.15)
    bordered_rect(sl, x, y, Inches(2.95), Inches(1.9), WHITE, color, Pt(2), 0.04)
    line_shape(sl, x + Inches(0.02), y + Inches(0.02), Inches(2.91), Inches(0.4), color)
    txt(sl, x + Inches(0.1), y + Inches(0.05), Inches(2.75), Inches(0.35),
        label, 11, WHITE, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), y + Inches(0.5), Inches(2.75), Inches(0.35),
        title, 14, DARK, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), y + Inches(0.9), Inches(2.75), Inches(0.9),
        desc, 11, GRAY)


# ============================================================
# SLIDE 8: System Architecture
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "System Architecture", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

layers = [
    ("React App", "(Frontend)", "Screens · Components\nVoice · Redux", BLUE),
    ("API Client", "(Axios)", "HTTP Requests\nJWT Interceptor", ORANGE),
    ("Django REST", "(Backend)", "Views · Serializers\nNLP · Whisper", GREEN),
    ("Database", "(PostgreSQL)", "10+ Models\nAudit Trail", LIGHT_BLUE),
]

for i, (title, sub, desc, color) in enumerate(layers):
    x = Inches(0.15 + i * 2.5)
    y = Inches(0.85)
    bordered_rect(sl, x, y, Inches(2.15), Inches(1.5), WHITE, color, Pt(2), 0.04)
    line_shape(sl, x + Inches(0.02), y + Inches(0.02), Inches(2.11), Inches(0.35), color)
    txt(sl, x + Inches(0.1), y + Inches(0.05), Inches(1.95), Inches(0.3),
        title, 14, WHITE, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), y + Inches(0.4), Inches(1.95), Inches(0.25),
        sub, 10, color, False, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), y + Inches(0.75), Inches(1.95), Inches(0.6),
        desc, 11, GRAY, False, PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, x + Inches(2.15), y + Inches(0.45), Inches(0.35), Inches(0.4),
            "↔", 18, DARK, True, PP_ALIGN.CENTER)

# Data Flow
txt(sl, Inches(0.3), Inches(2.6), Inches(9), Inches(0.35),
    "Data Flow", 18, DARK, True)

flow = [
    ("1", "User speaks\ninto mic"),
    ("2", "Web Speech API\ntranscribes live"),
    ("3", "Audio uploaded\nto backend"),
    ("4", "Whisper STT +\nnormalization"),
    ("5", "Intent & entity\nextraction"),
    ("6", "Confirm dialog\n→ execute"),
]

for i, (num, desc) in enumerate(flow):
    x = Inches(0.1 + i * 1.65)
    y = Inches(3.05)
    bordered_rect(sl, x, y, Inches(1.45), Inches(1.3), WHITE, BLUE, Pt(1), 0.04)
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), y + Inches(0.08), Inches(0.4), Inches(0.4))
    c.fill.solid()
    c.fill.fore_color.rgb = BLUE
    c.line.fill.background()
    ctf = c.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    txt(sl, x + Inches(0.05), y + Inches(0.55), Inches(1.35), Inches(0.65),
        desc, 10, GRAY, False, PP_ALIGN.CENTER)
    if i < 5:
        txt(sl, x + Inches(1.45), y + Inches(0.35), Inches(0.2), Inches(0.35),
            "→", 14, BLUE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: Voice Processing Pipeline
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Voice Processing Pipeline", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

# Example
bordered_rect(sl, Inches(0.3), Inches(0.8), Inches(9.4), Inches(1.2), CARD_BG, BLUE, Pt(1.5), 0.03)
txt(sl, Inches(0.5), Inches(0.85), Inches(3), Inches(0.25),
    "EXAMPLE COMMAND", 10, BLUE, True)
txt(sl, Inches(0.5), Inches(1.1), Inches(8.5), Inches(0.3),
    '"Update marks for roll 1, maths 95, hindi 88"', 16, DARK, True)
txt(sl, Inches(0.5), Inches(1.5), Inches(8.5), Inches(0.35),
    '→ Intent: UPDATE_MARKS  |  Entities: {roll: 1, marks: [{maths: 95}, {hindi: 88}]}', 12, GRAY)

# NLP features
txt(sl, Inches(0.3), Inches(2.2), Inches(9), Inches(0.3),
    "Intelligent NLP Features (4,000+ lines, no ML models)", 14, BLUE, True)

nlp = [
    ("Homophone Fix", '"won"→"1"\n"rule"→"roll"', BLUE),
    ("Subject Fix", '"mass"→"maths"\n"signs"→"science"', ORANGE),
    ("Number Split", '"12345678910" →\n"1,2,3,...,10"', GREEN),
    ("Range Expand", '"1 to 10" →\n"1,2,3,...,10"', GOLD),
    ("50+ Intents", "UPDATE_MARKS\nMARK_ATTENDANCE", LIGHT_BLUE),
]

for i, (title, desc, color) in enumerate(nlp):
    x = Inches(0.1 + i * 1.96)
    y = Inches(2.6)
    bordered_rect(sl, x, y, Inches(1.85), Inches(1.9), WHITE, color, Pt(1.5), 0.04)
    line_shape(sl, x + Inches(0.02), y + Inches(0.02), Inches(1.81), Inches(0.35), color)
    txt(sl, x + Inches(0.05), y + Inches(0.05), Inches(1.75), Inches(0.3),
        title, 12, WHITE, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.05), y + Inches(0.5), Inches(1.75), Inches(1.2),
        desc, 11, GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: App Features — Marks & Attendance
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "App Features — Marks & Attendance", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

features = [
    ("Subject Marks Entry", 'Voice: "Enter marks for roll 5, maths 92, hindi 88"'),
    ("Question-wise Marks", 'Voice: "Questions 1,2,3 marks as 4,5,3"'),
    ("Attendance Management", 'Voice: "Mark all present except roll 5, 12"'),
    ("Voice Navigation", 'Voice: "Open marks sheet class 5 section A"'),
    ("Auto Grade Calculation", "A+/A/B+/B/C/D/F calculated automatically from marks"),
    ("Student Reports", "View progress, marks summary & attendance percentage"),
]

feature_colors = [BLUE, ORANGE, GREEN, LIGHT_BLUE, GOLD, RGBColor(0xA5, 0xA5, 0xA5)]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.3 + col * 4.85)
    y = Inches(0.85 + row * 1.4)
    bordered_rect(sl, x, y, Inches(4.6), Inches(1.15), WHITE, feature_colors[i], Pt(1.5), 0.04)
    txt(sl, x + Inches(0.2), y + Inches(0.1), Inches(4.2), Inches(0.35),
        title, 16, feature_colors[i], True)
    txt(sl, x + Inches(0.2), y + Inches(0.5), Inches(4.2), Inches(0.55),
        desc, 12, GRAY)


# ============================================================
# SLIDE 11: Backend & Data Model
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Backend & Data Model", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

# Left: Models
models = [
    ("CustomUser", "Teacher / Admin roles"),
    ("ClassSection", "Class + Section + Year"),
    ("Student", "Roll, Name, DOB, Contact"),
    ("Marks / QuestionMarks", "Subject, Exam, Score, Grade"),
    ("Attendance", "Session + Record (P/A/L)"),
    ("VoiceCommand", "Audio, Intent, Entities, Status"),
]

model_colors = [BLUE, ORANGE, GREEN, LIGHT_BLUE, GOLD, RGBColor(0xA5, 0xA5, 0xA5)]

for i, (model, desc) in enumerate(models):
    y = Inches(0.8 + i * 0.65)
    bordered_rect(sl, Inches(0.3), y, Inches(4.5), Inches(0.55), WHITE, model_colors[i], Pt(1.5), 0.03)
    txt(sl, Inches(0.45), y + Inches(0.08), Inches(1.8), Inches(0.35),
        model, 12, model_colors[i], True)
    txt(sl, Inches(2.3), y + Inches(0.08), Inches(2.4), Inches(0.35),
        desc, 11, GRAY)

# Right: API endpoints
txt(sl, Inches(5.3), Inches(0.75), Inches(4.5), Inches(0.3),
    "REST API Endpoints", 16, BLUE, True)

endpoints = [
    "POST  /api/v1/auth/login/",
    "POST  /api/v1/auth/token/refresh/",
    "POST  /api/v1/voice/upload/",
    "POST  /api/v1/voice/commands/{id}/confirm/",
    "POST  /api/v1/voice/commands/{id}/reject/",
    "GET   /api/v1/marks/{class}/{section}/",
    "POST  /api/v1/marks/",
    "GET   /api/v1/attendance/{class}/{section}/",
    "POST  /api/v1/attendance/",
]

bordered_rect(sl, Inches(5.2), Inches(1.1), Inches(4.5), Inches(3.7), CARD_BG, BLUE, Pt(1), 0.03)
for i, ep in enumerate(endpoints):
    txt(sl, Inches(5.4), Inches(1.2 + i * 0.38), Inches(4.1), Inches(0.35),
        ep, 10, DARK)


# ============================================================
# SLIDE 12: Security & Edge Cases
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Security & Edge Cases", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

# Left: Security
txt(sl, Inches(0.5), Inches(0.8), Inches(4.5), Inches(0.35),
    "Security Measures", 18, BLUE, True)

security = [
    "JWT auth with auto token refresh (15min/7day)",
    "Role-Based Access Control per endpoint",
    "Intent filtering by user role on backend",
    "Two-phase commit: preview before DB write",
    "CORS headers · Django password hashing",
]
bordered_rect(sl, Inches(0.3), Inches(1.2), Inches(4.6), Inches(3.4), WHITE, BLUE, Pt(1.5), 0.03)
for i, s in enumerate(security):
    txt(sl, Inches(0.5), Inches(1.4 + i * 0.55), Inches(4.2), Inches(0.45),
        f"• {s}", 13, DARK)

# Right: Edge Cases
txt(sl, Inches(5.3), Inches(0.8), Inches(4.5), Inches(0.35),
    "Edge Cases Handled (20+)", 18, ORANGE, True)

edges = [
    "Merged number splitting (fast speech)",
    "Indian English accent normalization",
    "Incomplete command detection",
    "Student/class not found → friendly error",
    "Batch updates with mismatched lengths",
    "Decimal marks preservation (4.5)",
    "Duplicate attendance prevention",
]
bordered_rect(sl, Inches(5.1), Inches(1.2), Inches(4.6), Inches(3.4), WHITE, ORANGE, Pt(1.5), 0.03)
for i, e in enumerate(edges):
    txt(sl, Inches(5.3), Inches(1.4 + i * 0.45), Inches(4.2), Inches(0.4),
        f"• {e}", 13, DARK)


# ============================================================
# SLIDE 13: Advantages and Limitations
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Advantages and Limitations", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

# Left: Advantages
txt(sl, Inches(0.5), Inches(0.8), Inches(4.5), Inches(0.35),
    "Advantages", 20, GREEN, True)

advantages = [
    "Hands-free operation reduces admin workload by 30-40%",
    "Custom NLP works offline — no API costs or latency",
    "Two-phase commit prevents accidental data corruption",
    "Handles Indian English accents & speech patterns",
    "Batch operations: multiple marks in one command",
]
bordered_rect(sl, Inches(0.3), Inches(1.2), Inches(4.6), Inches(3.4), WHITE, GREEN, Pt(1.5), 0.03)
for i, a in enumerate(advantages):
    txt(sl, Inches(0.5), Inches(1.4 + i * 0.55), Inches(4.2), Inches(0.45),
        f"✅  {a}", 13, DARK)

# Right: Limitations
txt(sl, Inches(5.3), Inches(0.8), Inches(4.5), Inches(0.35),
    "Limitations", 20, ORANGE, True)

limitations = [
    "Requires active backend connection (not fully offline)",
    "Regex NLP may miss very unusual phrasing",
    "Whisper model needs ~150MB download on first use",
    "Browser mic permission required for voice input",
    "Currently English-only voice commands",
]
bordered_rect(sl, Inches(5.1), Inches(1.2), Inches(4.6), Inches(3.4), WHITE, ORANGE, Pt(1.5), 0.03)
for i, l in enumerate(limitations):
    txt(sl, Inches(5.3), Inches(1.4 + i * 0.55), Inches(4.2), Inches(0.45),
        f"⚠️  {l}", 13, DARK)


# ============================================================
# SLIDE 14: Future Scope
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "Future Scope", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

future = [
    "Fee Management module — voice-driven fee collection, defaulter tracking & reports",
    "LLM-based intent extraction — replace regex with GPT/Claude for natural understanding",
    "Multi-language support — Hindi, Marathi & regional language voice commands",
    "Mobile app — React Native with offline voice support",
    "Analytics dashboard — performance trends, charts & AI-powered insights",
]

for i, item in enumerate(future):
    y = Inches(0.85 + i * 0.8)
    txt(sl, Inches(0.5), y, Inches(9), Inches(0.6),
        f"• {item}", 16, DARK)


# ============================================================
# SLIDE 15: References
# ============================================================
sl = add_slide(layout_title_content)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
    "References", 28, DARK, True, PP_ALIGN.LEFT, "Calibri Light")

refs = [
    "OpenAI Whisper, Robust Speech Recognition via Large-Scale Weak Supervision, 2023",
    "Django Documentation, Django Software Foundation, djangoproject.com",
    "React 18 Documentation, Meta Platforms, react.dev",
    "Web Speech API Specification, W3C Community Group",
    "Django REST Framework, Tom Christie, django-rest-framework.org",
    "Redux Toolkit Documentation, Redux Team, redux-toolkit.js.org",
    "PostgreSQL 15 Documentation, PostgreSQL Global Development Group",
    "Tailwind CSS Documentation, Tailwind Labs, tailwindcss.com",
    "K12 Techno Services, k12technoservices.com",
]

for i, ref in enumerate(refs):
    txt(sl, Inches(0.5), Inches(0.8 + i * 0.45), Inches(9), Inches(0.4),
        f" {ref}", 13, DARK)


# ============================================================
# SLIDE 16: Thank You
# ============================================================
sl = add_slide(layout_section)
for ph in list(sl.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

txt(sl, Inches(0.5), Inches(1.5), Inches(9), Inches(1.0),
    "Thank You", 48, DARK, True, PP_ALIGN.CENTER, "Calibri Light")

txt(sl, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
    "B. Parshuram  |  Dr. Girdhar Gopal  |  ECE Dept.  |  NIT Patna", 16, GRAY, False, PP_ALIGN.CENTER)


# === SAVE ===
output = "/Users/bardawalparshuram/Desktop/doneATOA 2/Audio_to_Action/ReATOA_Viva_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
